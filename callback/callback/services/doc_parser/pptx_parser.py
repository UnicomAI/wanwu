"""PPTX 解析器：用 python-pptx 按 slide -> shape 提取文本，达到上限即停止。"""

from pptx import Presentation

from .base import DocumentParser


class PptxParser(DocumentParser):
    def parse(self, file_path: str, max_tokens: int = 0) -> str:
        limit = self.max_chars(max_tokens)
        parts = []
        total = 0
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                text = self._shape_text(shape)
                if text:
                    parts.append(text)
                    parts.append("\n")
                    total += len(text) + 1
                    if limit is not None and total >= limit:
                        return self.truncate("".join(parts), max_tokens)
        return self.truncate("".join(parts), max_tokens)

    @staticmethod
    def _shape_text(shape) -> str:
        """提取 shape 文本：文本框直接取 text；表格按行列拼接。"""
        if shape.has_text_frame:
            return shape.text_frame.text
        if shape.has_table:
            lines = []
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                lines.append("\t".join(cells))
            return "\n".join(lines)
        return ""
