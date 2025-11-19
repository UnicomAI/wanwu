from pptx import Presentation
from typing import List, Optional

from langchain_core.documents import Document
from langchain_community.document_loaders import TextLoader

import os
import nltk
current_file_path = os.path.abspath(__file__)
# Get the directory where the current file is located
current_dir = os.path.dirname(current_file_path)
# Add project root directory to sys.path
# sys.path.append(root_dir)
# Splice the path to the nltk_data folder
nltk_data_path = os.path.join(current_dir, 'nltk_data')
nltk.data.path.append(nltk_data_path)
# import minio_utils
from logging_config import setup_logging
logger_name = 'rag_pptx_loader'
app_name = os.getenv("LOG_FILE")
logger = setup_logging(app_name,logger_name)
logger.info(logger_name+'---------LOG_FILE：'+repr(app_name))


def table_convert_html(table):
    # embedding_content = []
    #
    # table_title_list = []

    def is_empty(cell):
        return cell is None
        # return cell is None or str(cell).strip() == ''

    def get_colspan(row, col_idx, processed_cells, row_idx):
        # Calculate the number of columns spanned by the current cell
        if (row_idx, col_idx) in processed_cells:
            return 0  # Skip processed cells
        # Calculate the number of columns spanned by the current cell
        colspan = 1
        for i in range(col_idx + 1, len(row)):
            if is_empty(row[i]) and (row_idx, i) not in processed_cells:
                processed_cells.add((row_idx, i))  # Flag processed
                colspan += 1
            else:
                break
        return colspan

    def get_rowspan(rows, row_idx, col_idx, processed_cells):
        if (row_idx, col_idx) in processed_cells:
            return 0  # Skip processed cells
        rowspan = 1
        for i in range(row_idx + 1, len(rows)):
            if is_empty(rows[i][col_idx]) and (i, col_idx) not in processed_cells:
                processed_cells.add((i, col_idx))  # Flag processed
                rowspan += 1
            else:
                break
        return rowspan

    # Start building an HTML table
    html = "<table border='1'>\n"
    processed_cells = set()  # Record the processed cell position (row_idx, col_idx)
    for row_idx, row in enumerate(table):
        html += "<tr>"
        row_text = "<tr>"
        col_idx = 0
        while col_idx < len(row):
            if (row_idx, col_idx) in processed_cells:
                col_idx += 1  # Skip processed cells to avoid infinite loops
                continue
            cell = row[col_idx]

            if is_empty(cell):
                col_idx += 1
                continue  # Skip empty cells

            # Get the number of spanned columns
            colspan = get_colspan(row, col_idx, processed_cells, row_idx)

            # Get the number of spanned rows
            rowspan = get_rowspan(table, row_idx, col_idx, processed_cells)

            # Add cell
            if colspan > 1 or rowspan > 1:
                html += f"<td colspan='{colspan}' rowspan='{rowspan}'>{cell}</td>"
                row_text += f"<td colspan='{colspan}' rowspan='{rowspan}'>{cell}</td>"
            else:
                html += f"<td>{cell}</td>"
                row_text += f"<td>{cell}</td>"
            # The contents of the first two columns except the first row are called embedding indexes
            # if row_idx > 0 and col_idx < 2 and str(cell).strip() != '':
            #     # For tables without row dividing lines in the table: cut the index by newline character
            #     if len(table) <= 2 and "\n" in str(cell):
            #         cell_list = str(cell).split("\n")
            #         embedding_content.extend(cell_list)
            #     else:
            #         embedding_content.append(cell)
            # Mark cells that have been processed
            for i in range(rowspan):
                for j in range(colspan):
                    processed_cells.add((row_idx + i, col_idx + j))
            # Skip already processed cross-column cells
            col_idx += colspan

        html += "</tr>\n"
        row_text += "</tr>"
        # if row_idx > 0:
        #     embedding_content.append(row_text)

    # End HTML table
    html += "</table>\n"
    # print(json.dumps(embedding_content,ensure_ascii=False))
    return html


class PPTXLoader(TextLoader):
    def load(self) -> List[Document]:
        text = ""
        try:
            prs = Presentation(self.file_path)
            print(prs)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        t = text_frame.text
                        text += t + '\n'
                    elif shape.has_table:
                        one_table_data = []
                        for row in shape.table.rows:  # read each line
                            row_data = []
                            for cell in row.cells:  # Read all cells in a row
                                if cell.text != "":
                                    row_data.append(cell.text)
                                else:
                                    row_data.append(None)
                                # cell.text = cell.text if cell.text != "" else ""
                                # c = cell.text
                                # row_data.append(c)
                            one_table_data.append(row_data)  # Store each row in the table

                        print("one_table_data=%s" % one_table_data)
                        table_html = table_convert_html(one_table_data)
                        text += table_html + '\n'
        except Exception as e:
            raise RuntimeError(f"Error loading {self.file_path}") from e

        metadata = {"source": self.file_path}
        return [Document(page_content=text, metadata=metadata)]
if __name__ == "__main__":

    filepath = "./your_file.pptx"
    loader = PPTXLoader(filepath)
    docs = loader.load()
    for doc in docs:
        print(doc)